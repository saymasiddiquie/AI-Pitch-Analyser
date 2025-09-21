import streamlit as st
# Try to import a compatible PDF reader implementation. Do not crash if unavailable.
try:
    from PyPDF2 import PdfReader as _PdfReader
except ImportError:
    try:
        from pypdf import PdfReader as _PdfReader
    except ImportError:
        _PdfReader = None
# Try to import PowerPoint reader. Optional dependency in some environments.
try:
    from pptx import Presentation as _Presentation
except ImportError:
    _Presentation = None
import io
import re
import os
# Load .env if python-dotenv is available; otherwise define a no-op.
try:
    from dotenv import load_dotenv
except ImportError:
    def load_dotenv(*args, **kwargs):
        return False
# Try to import Google Generative AI SDK; don't crash if missing
try:
    import google.generativeai as genai
except ImportError:
    genai = None

# Remove or comment out print statements in production
# print("Module imported successfully!")

# ---- INIT GEMINI ----
load_dotenv()
GEMINI_API_KEY = None
try:
    # Prefer Streamlit Secrets if available
    GEMINI_API_KEY = st.secrets.get("GEMINI_API_KEY")
except Exception:
    pass
if not GEMINI_API_KEY:
    GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if genai is not None and GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel("gemini-1.5-flash")
    except Exception:
        model = None
else:
    model = None

# ---- PAGE CONFIG ----
st.set_page_config(
    page_title="AI Pitch Deck Summariser",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# Hide the hamburger menu and footer for a cleaner look
hide_streamlit_style = """
<style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    .stDeployButton {visibility: hidden;}
</style>
"""
st.markdown(hide_streamlit_style, unsafe_allow_html=True)

# Hide the "Made with Streamlit" footer
hide_streamlit_footer = """
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    </style>
"""
st.markdown(hide_streamlit_footer, unsafe_allow_html=True)

# ---- BEAUTIFUL CUSTOM CSS ----
st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap');
        
        .stApp {
            background: linear-gradient(135deg, #f5f0ff 0%, #e6e6ff 100%);
            font-family: 'Inter', sans-serif;
            color: #2d2d5a;
        }
        
        .main-header {
            background: linear-gradient(135deg, #9c89ff 0%, #6a5acd 100%);
            padding: 4rem 0;
            text-align: center;
            border-radius: 20px;
            margin: 2rem 0 3rem 0;
            box-shadow: 0 15px 35px rgba(156, 137, 255, 0.2);
            border: 1px solid rgba(255, 255, 255, 0.3);
            position: relative;
            overflow: hidden;
        }
        
        .main-header::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            bottom: 0;
            background: radial-gradient(circle at 20% 30%, rgba(255,255,255,0.2) 0%, transparent 60%);
            pointer-events: none;
        }
        
        .main-title {
            font-size: 4rem;
            font-weight: 800;
            background: linear-gradient(45deg, #ffffff, #f0e6ff);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            margin-bottom: 1rem;
            text-shadow: 0 2px 4px rgba(0,0,0,0.1);
            letter-spacing: -0.5px;
            position: relative;
            z-index: 1;
        }
        
        .subtitle {
            color: #f0f0ff;
            font-size: 1.4rem;
            font-weight: 400;
            margin: 0.5rem 0 2rem 0;
            max-width: 800px;
            line-height: 1.6;
            position: relative;
            z-index: 1;
            text-shadow: 0 1px 2px rgba(0,0,0,0.1);
        }
        
        .upload-container {
            background: rgba(255, 255, 255, 0.98);
            backdrop-filter: blur(20px);
            padding: 2.5rem;
            border-radius: 24px;
            box-shadow: 0 15px 35px rgba(156, 137, 255, 0.15);
            border: 2px dashed #b8a9ff;
            margin: 2rem 0 3rem 0;
            transition: all 0.3s ease;
            text-align: center;
        }
        
        .upload-container:hover {
            transform: translateY(-3px);
            box-shadow: 0 20px 45px rgba(156, 137, 255, 0.25);
            border-color: #9c89ff;
        }
        
        .stButton > button {
            background: linear-gradient(135deg, #8a7aff 0%, #6a5acd 100%);
            color: white;
            border: none;
            padding: 1.1rem 3.5rem;
            border-radius: 50px;
            font-size: 1.2rem;
            font-weight: 600;
            box-shadow: 0 10px 30px rgba(138, 122, 255, 0.4);
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            width: auto;
            margin: 1.5rem auto 0;
            display: inline-block;
            min-width: 250px;
            text-transform: uppercase;
            letter-spacing: 0.5px;
        }
        
        .stButton > button:hover {
            transform: translateY(-3px) scale(1.02);
            box-shadow: 0 20px 45px rgba(138, 122, 255, 0.5);
            background: linear-gradient(135deg, #6a5acd 0%, #8a7aff 100%);
        }
        
        .analysis-card {
            background: rgba(255, 255, 255, 0.95);
            backdrop-filter: blur(20px);
            padding: 2rem;
            border-radius: 20px;
            box-shadow: 0 20px 40px rgba(0,0,0,0.15);
            border: 1px solid rgba(255,255,255,0.2);
            margin: 2rem 0;
            animation: slideUp 0.8s ease-out;
        }
        
        @keyframes slideUp {
            from {
                opacity: 0;
                transform: translateY(30px);
            }
            to {
                opacity: 1;
                transform: translateY(0);
            }
        }
        
        .section-header {
            background: linear-gradient(135deg, #667eea, #764ba2);
            color: white;
            padding: 1rem 1.5rem;
            border-radius: 15px;
            font-size: 1.3rem;
            font-weight: 600;
            margin: 1.5rem 0 1rem 0;
            box-shadow: 0 8px 25px rgba(102, 126, 234, 0.3);
            display: flex;
            align-items: center;
            gap: 10px;
        }
        
        .bullet-container {
            background: linear-gradient(135deg, #f9f7ff, #f0ecff);
            padding: 1.5rem;
            border-radius: 15px;
            margin: 1rem 0;
            border-left: 5px solid #8a7aff;
            box-shadow: 0 5px 20px rgba(156, 137, 255, 0.1);
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            position: relative;
            overflow: hidden;
        }
        
        .bullet-container::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: linear-gradient(135deg, rgba(138, 122, 255, 0.1), rgba(106, 90, 205, 0.1));
            opacity: 0;
            transition: opacity 0.3s ease;
        }
        
        .bullet-container:hover::before {
            opacity: 1;
        }
        
        .bullet-container:hover {
            transform: translateX(5px);
        }
        
        .bullet-point {
            color: #3d3d6b;
            font-size: 1.1rem;
            font-weight: 500;
            line-height: 1.7;
            margin: 0;
            position: relative;
            padding-left: 1.5rem;
        }
        
        .bullet-point::before {
            content: '•';
            color: #8a7aff;
            font-size: 1.8rem;
            position: absolute;
            left: -0.2rem;
            top: -0.4rem;
            line-height: 1;
        }
        
        .highlight-metric {
            background: linear-gradient(135deg, #8a7aff, #6a5acd);
            color: white;
            padding: 1rem 2rem;
            border-radius: 30px;
            display: inline-block;
            font-weight: 600;
            margin: 0.5rem;
            box-shadow: 0 10px 25px rgba(138, 122, 255, 0.3);
            font-size: 1.1rem;
            border: none;
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            text-shadow: 0 1px 2px rgba(0,0,0,0.1);
        }
        
        .highlight-metric:hover {
            transform: translateY(-2px);
            box-shadow: 0 15px 30px rgba(138, 122, 255, 0.4);
        }
        
        .investment-card {
            background: linear-gradient(135deg, #667eea, #764ba2);
            color: white;
            padding: 1.5rem;
            border-radius: 20px;
            margin: 1rem 0;
            box-shadow: 0 15px 35px rgba(102, 126, 234, 0.4);
        }
        
        .competitor-badge {
            background: rgba(138, 122, 255, 0.1);
            border: 2px solid rgba(138, 122, 255, 0.3);
            padding: 1rem 1.5rem;
            border-radius: 15px;
            margin: 0.8rem 0;
            color: #4a3b8a;
            font-weight: 600;
            font-size: 1.05rem;
            transition: all 0.3s ease;
            display: inline-block;
            width: 100%;
        }
        
        .competitor-badge:hover {
            background: rgba(138, 122, 255, 0.15);
            border-color: #8a7aff;
            transform: translateX(5px);
        }
        
        .loading-text {
            text-align: center;
            color: #667eea;
            font-size: 1.2rem;
            font-weight: 600;
            margin: 2rem 0;
        }
        
        .download-btn {
            background: linear-gradient(135deg, #8a7aff, #6a5acd);
            color: white;
            padding: 1rem 2.5rem;
            border-radius: 50px;
            font-weight: 600;
            font-size: 1.1rem;
            box-shadow: 0 10px 25px rgba(138, 122, 255, 0.3);
            margin: 2rem auto 1rem;
            display: inline-block;
            text-decoration: none;
            transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
            border: none;
            cursor: pointer;
            text-align: center;
            min-width: 250px;
        }
        
        .download-btn:hover {
            transform: translateY(-2px) scale(1.02);
            box-shadow: 0 15px 35px rgba(138, 122, 255, 0.4);
            color: white;
        }
        
        .header-content {
            display: flex;
            flex-direction: column;
            align-items: center;
            text-align: center;
            padding: 2rem;
        }
        
        .header-illustration {
            display: flex;
            justify-content: center;
            padding: 2rem;
        }
        
        .header-illustration img {
            width: 200px;
            height: 200px;
            border-radius: 50%;
        }
        
        .cta-buttons {
            display: flex;
            gap: 1rem;
            margin-top: 1rem;
        }
        
        .cta-button {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 0.8rem 2rem;
            border-radius: 50px;
            font-weight: 600;
            box-shadow: 0 8px 20px rgba(102, 126, 234, 0.3);
            text-decoration: none;
        }
        
        .cta-button:hover {
            transform: translateY(-3px);
            box-shadow: 0 15px 40px rgba(102, 126, 234, 0.6);
            background: linear-gradient(135deg, #764ba2 0%, #667eea 100%);
        }
        
        .primary {
            background: linear-gradient(135deg, #10b981, #059669);
            color: white;
            padding: 0.8rem 2rem;
            border-radius: 50px;
            font-weight: 600;
            box-shadow: 0 8px 20px rgba(16, 185, 129, 0.3);
        }
        
        .primary:hover {
            transform: translateY(-3px);
            box-shadow: 0 15px 40px rgba(16, 185, 129, 0.6);
            background: linear-gradient(135deg, #059669, #10b981);
        }
        
        .secondary {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 0.8rem 2rem;
            border-radius: 50px;
            font-weight: 600;
            box-shadow: 0 8px 20px rgba(102, 126, 234, 0.3);
        }
        
        .secondary:hover {
            transform: translateY(-3px);
            box-shadow: 0 15px 40px rgba(102, 126, 234, 0.6);
            background: linear-gradient(135deg, #764ba2 0%, #667eea 100%);
        }
        
        .features-section {
            padding: 2rem;
            text-align: center;
        }
        
        .features-grid {
            display: grid;
            grid-template-columns: repeat(3, 1fr);
            gap: 2rem;
        }
        
        .feature-card {
            background: rgba(255, 255, 255, 0.98);
            backdrop-filter: blur(20px);
            padding: 2.5rem 2rem;
            border-radius: 20px;
            box-shadow: 0 15px 35px rgba(156, 137, 255, 0.1);
            border: 1px solid rgba(255,255,255,0.8);
            transition: all 0.3s ease;
            text-align: center;
            position: relative;
            overflow: hidden;
            z-index: 1;
        }
        
        .feature-card::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: linear-gradient(90deg, #8a7aff, #6a5acd);
            transition: all 0.3s ease;
            opacity: 0;
        }
        
        .feature-card:hover {
            transform: translateY(-5px);
            box-shadow: 0 20px 45px rgba(156, 137, 255, 0.2);
        }
        
        .feature-card:hover::before {
            opacity: 1;
        }
        
        .feature-icon {
            font-size: 2.5rem;
            margin-bottom: 1.5rem;
            display: inline-block;
            background: linear-gradient(135deg, #8a7aff 0%, #6a5acd 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            background-clip: text;
            color: transparent;
        }
        
        .section-spacer {
            margin-top: 4rem;
        }
        
        .steps-container {
            display: flex;
            flex-direction: column;
            align-items: center;
            padding: 2rem;
        }
        
        .step {
            display: flex;
            align-items: center;
            gap: 1rem;
            margin-bottom: 2rem;
        }
        
        .step-number {
            font-size: 2.5rem;
            font-weight: 700;
            margin-right: 1.5rem;
            color: #8a7aff;
            min-width: 50px;
            height: 50px;
            background: rgba(138, 122, 255, 0.1);
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            box-shadow: 0 5px 15px rgba(138, 122, 255, 0.2);
            position: relative;
            z-index: 1;
        }
        
        .step-number::after {
            content: '';
            position: absolute;
            top: -3px;
            left: -3px;
            right: -3px;
            bottom: -3px;
            border: 2px solid rgba(138, 122, 255, 0.3);
            border-radius: 50%;
            z-index: -1;
            animation: pulse 2s infinite;
        }
        
        @keyframes pulse {
            0% { transform: scale(1); opacity: 0.7; }
            50% { transform: scale(1.1); opacity: 0.4; }
            100% { transform: scale(1); opacity: 0.7; }
        }
        
        .step-content {
            padding: 1rem;
            border-radius: 15px;
            background: rgba(255, 255, 255, 0.95);
            backdrop-filter: blur(20px);
            box-shadow: 0 10px 30px rgba(0,0,0,0.1);
        }
        
        .section-title {
            font-size: 2.5rem;
            color: #4a3b8a;
            margin: 4rem 0 2rem 0;
            text-align: center;
            font-weight: 700;
            position: relative;
            display: inline-block;
            left: 50%;
            transform: translateX(-50%);
            padding: 0 2rem;
        }
        
        .section-title::after {
            content: '';
            position: absolute;
            bottom: -10px;
            left: 50%;
            transform: translateX(-50%);
            width: 80px;
            height: 4px;
            background: linear-gradient(90deg, #8a7aff, #6a5acd);
            border-radius: 2px;
        }
        
        .step-content h4 {
            color: #4a3b8a;
            margin-top: 0;
            margin-bottom: 0.5rem;
            font-size: 1.3rem;
        }
        
        .step-content p {
            color: #6b6b8a;
            margin: 0;
            font-size: 1rem;
            line-height: 1.6;
        }
    </style>
""", unsafe_allow_html=True)

# ---- BEAUTIFUL HEADER ----
st.markdown("""
    <div class="main-header">
        <div class="header-content">
            <h1 class="main-title">🚀 AI-Powered Pitch Deck Analysis</h1>
            <p class="subtitle">Get instant, comprehensive insights into any pitch deck with our AI-powered analysis platform</p>
            <div class="cta-buttons">
                <a href="#upload-section" class="cta-button primary">Analyze Now</a>
                <a href="#features" class="cta-button secondary">Learn More</a>
            </div>
        </div>
        <div class="header-illustration">
            <img src="https://img.icons8.com/clouds/300/000000/rocket.png" alt="Rocket Illustration">
        </div>
    </div>
""", unsafe_allow_html=True)

# Features Section
st.markdown("<div id='features'></div>", unsafe_allow_html=True)
st.markdown("""
    <div class="features-section">
        <h2 class="section-title">Why Choose Our Pitch Deck Analyzer?</h2>
        <div class="features-grid">
            <div class="feature-card">
                <div class="feature-icon">⚡</div>
                <h3>Lightning Fast</h3>
                <p>Get detailed analysis in seconds, not hours</p>
            </div>
            <div class="feature-card">
                <div class="feature-icon">🎯</div>
                <h3>Actionable Insights</h3>
                <p>Clear, practical recommendations to improve your pitch</p>
            </div>
            <div class="feature-card">
                <div class="feature-icon">🔍</div>
                <h3>Deep Analysis</h3>
                <p>Comprehensive evaluation of all key aspects</p>
            </div>
        </div>
    </div>
""", unsafe_allow_html=True)

# Upload Section
st.markdown("<div id='upload-section' class='section-spacer'></div>", unsafe_allow_html=True)
st.markdown("<h2 class='section-title'>Upload Your Pitch Deck</h2>", unsafe_allow_html=True)

# How It Works Section
with st.expander("ℹ️ How It Works", expanded=False):
    st.markdown("""
    <div class="steps-container">
        <div class="step">
            <div class="step-number">1</div>
            <div class="step-content">
                <h4>Upload Your Pitch Deck</h4>
                <p>Upload a PDF or PowerPoint file (PPTX) of your pitch deck</p>
            </div>
        </div>
        <div class="step">
            <div class="step-number">2</div>
            <div class="step-content">
                <h4>AI Analysis</h4>
                <p>Our AI will analyze your deck's content, structure, and messaging</p>
            </div>
        </div>
        <div class="step">
            <div class="step-number">3</div>
            <div class="step-content">
                <h4>Get Insights</h4>
                <p>Receive detailed feedback and recommendations</p>
            </div>
        </div>
    </div>
""", unsafe_allow_html=True)

st.markdown("### 📁 Upload Your Pitch Deck")
uploaded_file = st.file_uploader(
    "Upload your pitch deck",
    type=["pdf", "pptx"],
    accept_multiple_files=False,
    help="Supported formats: PDF, PPTX",
    label_visibility="collapsed"
)
st.markdown('</div>', unsafe_allow_html=True)

# ---- ANALYSIS SECTION ----
if uploaded_file:
    col1, col2, col3 = st.columns([1,2,1])
    with col2:
        if st.button("🔍 Generate AI Analysis"):
            text = ""
            
            # File extraction with progress
            with st.spinner("📖 Extracting content from your pitch deck..."):
                # PDF extraction
                if uploaded_file.name.endswith(".pdf"):
                    if _PdfReader is None:
                        st.error("PDF reader dependency not found. Please ensure either PyPDF2 or pypdf is installed.")
                    else:
                        reader = _PdfReader(uploaded_file)
                        for page in reader.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                
                # PPTX extraction
                elif uploaded_file.name.endswith(".pptx"):
                    if _Presentation is None:
                        st.error("PPTX support is unavailable. Please ensure the 'python-pptx' package is installed.")
                    else:
                        prs = _Presentation(io.BytesIO(uploaded_file.read()))
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    text += shape.text + "\n"
            
            # ---- AI ANALYSIS ----
            if text:
                with st.spinner("🤖 AI is analyzing your pitch deck..."):
                    prompt = f"""
                    You are a senior venture capitalist and business analyst. Analyze this pitch deck and provide precise, actionable insights in the following format:

                    **EXECUTIVE SUMMARY**
                    • Company: [One line company description]
                    • Problem: [Problem they solve in one line]  
                    • Solution: [Solution in one line]
                    • Market: [Target market size/opportunity in one line]

                    **STRENGTHS**
                    • [Strength 1 - max 12 words]
                    • [Strength 2 - max 12 words]
                    • [Strength 3 - max 12 words]

                    **WEAKNESSES**
                    • [Weakness 1 - max 12 words]
                    • [Weakness 2 - max 12 words]
                    • [Weakness 3 - max 12 words]

                    **OPPORTUNITIES**
                    • [Opportunity 1 - max 12 words]
                    • [Opportunity 2 - max 12 words]
                    • [Opportunity 3 - max 12 words]

                    **THREATS**
                    • [Threat 1 - max 12 words]
                    • [Threat 2 - max 12 words]
                    • [Threat 3 - max 12 words]

                    **KEY RISKS**
                    • [Risk 1 - max 15 words]
                    • [Risk 2 - max 15 words]
                    • [Risk 3 - max 15 words]

                    **MAIN COMPETITORS**
                    • [Competitor 1]: [Why they compete - max 10 words]
                    • [Competitor 2]: [Why they compete - max 10 words]
                    • [Competitor 3]: [Why they compete - max 10 words]

                    **INVESTMENT ANALYSIS**
                    • Stage: [Pre-seed/Seed/Series A/B/C]
                    • Range: $[X]M - $[Y]M
                    • Primary Use: [Main use of funds - max 10 words]
                    • Secondary Use: [Second use of funds - max 10 words]
                    • Timeline: [Expected fundraising timeline - max 8 words]

                    Be precise, direct, and actionable. Each bullet point should provide clear value.

                    Pitch Deck Content:
                    {text}
                    """

                    if model is None:
                        st.error("AI model is not available. Ensure 'google-generativeai' is installed and GEMINI_API_KEY is set in Streamlit Secrets or environment variables.")
                        analysis = None
                    else:
                        try:
                            response = model.generate_content(prompt)
                            analysis = getattr(response, 'text', '') or ''
                        except Exception as e:
                            st.error(f"AI request failed: {e}")
                            analysis = None

                # ---- DISPLAY BEAUTIFUL ANALYSIS ----
                if analysis:
                    st.markdown('<div class="analysis-card">', unsafe_allow_html=True)
                    st.markdown("## 📊 AI Analysis Results")
                
                # Parse and display with beautiful formatting
                sections = analysis.split("**")
                icons = {
                    "EXECUTIVE SUMMARY": "🎯",
                    "STRENGTHS": "💪",
                    "WEAKNESSES": "⚠️",
                    "OPPORTUNITIES": "🚀",
                    "THREATS": "⚡",
                    "KEY RISKS": "🚨",
                    "MAIN COMPETITORS": "🏆",
                    "INVESTMENT ANALYSIS": "💰"
                }
                
                for i in range(1, len(sections), 2):
                    if i < len(sections):
                        section_title = sections[i].strip()
                        section_content = sections[i+1].strip() if i+1 < len(sections) else ""
                        
                        if section_title in icons:
                            st.markdown(f'<div class="section-header">{icons[section_title]} {section_title}</div>', unsafe_allow_html=True)
                            
                            # Special formatting for investment section
                            if section_title == "INVESTMENT ANALYSIS":
                                st.markdown('<div class="investment-card">', unsafe_allow_html=True)
                                for line in section_content.split('\n'):
                                    if line.strip() and line.strip().startswith('•'):
                                        point = line.strip().replace('•', '').strip()
                                        st.markdown(f'<div class="highlight-metric">{point}</div>', unsafe_allow_html=True)
                                st.markdown('</div>', unsafe_allow_html=True)
                            
                            # Special formatting for competitors
                            elif section_title == "MAIN COMPETITORS":
                                for line in section_content.split('\n'):
                                    if line.strip() and line.strip().startswith('•'):
                                        point = line.strip().replace('•', '').strip()
                                        st.markdown(f'<div class="competitor-badge">🏢 {point}</div>', unsafe_allow_html=True)
                            
                            # Regular bullet points for other sections
                            else:
                                for line in section_content.split('\n'):
                                    if line.strip() and line.strip().startswith('•'):
                                        point = line.strip().replace('•', '').strip()
                                        st.markdown(f'<div class="bullet-container"><p class="bullet-point">• {point}</p></div>', unsafe_allow_html=True)

                    st.markdown('</div>', unsafe_allow_html=True)
                    
                    # Download button
                    col1, col2, col3 = st.columns([1,1,1])
                    with col2:
                        st.download_button(
                            label="📄 Download Analysis Report",
                            data=analysis,
                            file_name=f"pitch_analysis_{uploaded_file.name.split('.')[0]}.txt",
                            mime="text/plain",
                            help="Download the complete analysis as a text file"
                        )
            
            else:
                st.error("❌ Could not extract text from the uploaded file. Please ensure it contains readable content.")

# ---- FOOTER ----
st.markdown("---")
st.markdown(
    "<p style='text-align: center; color: rgba(255,255,255,0.7); font-size: 0.9rem;'>✨ Powered by Google Gemini AI | Built with Streamlit</p>", 
    unsafe_allow_html=True
)
